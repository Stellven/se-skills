"""Generate PPTX decks from documents with the Hermes PowerPoint skill.

This module provides two execution paths:

1. `--agent` launches a fresh `codex exec` run and asks Codex to use the
   NousResearch/Hermes PowerPoint skill installed in CODEX_HOME.
2. The default local path creates a deterministic PPTX directly with
   python-pptx. It is useful for smoke tests and environments where a nested
   Codex run is not desired.
"""

from __future__ import annotations

import argparse
from dataclasses import dataclass
import os
from pathlib import Path
import re
import shutil
import subprocess
import sys
from textwrap import shorten

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError as exc:  # pragma: no cover - surfaced in main() with a clear error.
    Presentation = None
    PPTX_IMPORT_ERROR = exc
else:
    PPTX_IMPORT_ERROR = None

try:
    from docx import Document as WordDocument
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph
except ImportError as exc:  # pragma: no cover - surfaced only for Word inputs.
    WordDocument = None
    DOCX_IMPORT_ERROR = exc
else:
    DOCX_IMPORT_ERROR = None

try:
    from pypdf import PdfReader
except ImportError as exc:  # pragma: no cover - surfaced only for PDF inputs.
    PdfReader = None
    PDF_IMPORT_ERROR = exc
else:
    PDF_IMPORT_ERROR = None

try:
    import fitz
except ImportError as exc:  # pragma: no cover - surfaced only for PDF figure extraction.
    fitz = None
    FITZ_IMPORT_ERROR = exc
else:
    FITZ_IMPORT_ERROR = None


HERMES_REPO = "NousResearch/hermes-agent"
HERMES_SKILL_PATH = "skills/productivity/powerpoint"
HERMES_SKILL_DIRNAME = "hermes-powerpoint"
WORD_EXTENSIONS = {".docx", ".docm"}
COMMON_DOCUMENT_HEADINGS = {
    "abstract",
    "introduction",
    "background",
    "method",
    "methods",
    "methodology",
    "results",
    "discussion",
    "conclusion",
    "conclusions",
    "references",
    "appendix",
}


@dataclass
class Section:
    title: str
    bullets: list[str]


@dataclass
class FigureAsset:
    label: str
    target_section: str
    caption: str
    path: Path


PDF_FIGURE_SPECS = (
    {
        "label": "Figure 1",
        "target_section": "1 Introduction",
        "caption": "Figure 1: Magentic-One completing a complex GAIA task",
        "top": 278,
    },
    {
        "label": "Figure 2",
        "target_section": "4 Magentic-One Overview",
        "caption": "Figure 2: Orchestrator-led multi-agent workflow",
        "top": 72,
    },
    {
        "label": "Figure 3",
        "target_section": "5.3 Ablations",
        "caption": "Figure 3: Ablation performance by level and capability",
        "top": 72,
    },
)


THEMES = {
    "executive": {
        "background": (246, 248, 251),
        "panel": (24, 36, 56),
        "accent": (0, 168, 150),
        "accent_2": (255, 192, 88),
        "text": (30, 41, 59),
        "muted": (100, 116, 139),
        "light": (255, 255, 255),
    },
    "charcoal": {
        "background": (245, 245, 242),
        "panel": (54, 69, 79),
        "accent": (184, 80, 66),
        "accent_2": (167, 190, 174),
        "text": (36, 42, 46),
        "muted": (96, 105, 112),
        "light": (255, 255, 255),
    },
    "teal": {
        "background": (241, 250, 248),
        "panel": (2, 80, 89),
        "accent": (2, 195, 154),
        "accent_2": (247, 205, 94),
        "text": (23, 48, 53),
        "muted": (82, 107, 111),
        "light": (255, 255, 255),
    },
}


def codex_home() -> Path:
    return Path(os.environ.get("CODEX_HOME", Path.home() / ".codex")).expanduser()


def hermes_skill_dir() -> Path:
    return codex_home() / "skills" / HERMES_SKILL_DIRNAME


def resolve_theme(theme_name: str) -> dict[str, RGBColor]:
    if Presentation is None:
        raise SystemExit(
            "python-pptx is required. Install dependencies with "
            "`pip install -r requirements.txt`. "
            f"Original import error: {PPTX_IMPORT_ERROR}"
        )
    return {key: RGBColor(*value) for key, value in THEMES[theme_name].items()}


def install_hermes_skill() -> None:
    """Install the upstream Hermes PowerPoint skill into CODEX_HOME."""
    destination = hermes_skill_dir()
    if (destination / "SKILL.md").exists():
        print(f"Hermes PowerPoint skill is already installed at {destination}")
        return

    installer = (
        codex_home()
        / "skills"
        / ".system"
        / "skill-installer"
        / "scripts"
        / "install-skill-from-github.py"
    )
    if not installer.exists():
        raise SystemExit(
            "Codex skill-installer was not found. Install it first or run this "
            "from a Codex environment with system skills available."
        )

    command = [
        sys.executable,
        str(installer),
        "--repo",
        HERMES_REPO,
        "--path",
        HERMES_SKILL_PATH,
        "--name",
        HERMES_SKILL_DIRNAME,
    ]
    subprocess.run(command, check=True)
    print("Restart Codex to pick up the newly installed skill.")


def read_document(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in WORD_EXTENSIONS:
        return read_word_document(path)
    if suffix == ".doc":
        return read_legacy_word_document(path)
    if suffix == ".pdf":
        return read_pdf_document(path)

    raw = path.read_text(encoding="utf-8", errors="ignore")
    if suffix == ".rtf":
        return read_rtf_document(path, raw)
    return raw


def read_word_document(path: Path) -> str:
    """Extract deck-ready text from a modern Word document."""
    if WordDocument is None:
        raise SystemExit(
            "python-docx is required for Word inputs. Install dependencies with "
            "`pip install -r requirements.txt`. "
            f"Original import error: {DOCX_IMPORT_ERROR}"
        )

    document = WordDocument(path)
    lines: list[str] = []
    table_index = 1
    for block in iter_word_blocks(document):
        if isinstance(block, Paragraph):
            converted = word_paragraph_to_markdown(block)
            if converted:
                lines.append(converted)
        elif isinstance(block, Table):
            table_lines = word_table_to_markdown(block, table_index)
            if table_lines:
                lines.extend(table_lines)
                table_index += 1

    return "\n\n".join(lines)


def iter_word_blocks(document) -> object:
    """Yield paragraphs and tables in document order."""
    for child in document.element.body.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, document)
        elif isinstance(child, CT_Tbl):
            yield Table(child, document)


def word_paragraph_to_markdown(paragraph: Paragraph) -> str | None:
    text = clean_inline_markdown(paragraph.text)
    if not text:
        return None

    style_name = (paragraph.style.name if paragraph.style else "").strip()
    style_lower = style_name.lower()
    if style_lower == "title":
        return f"# {text}"

    heading_match = re.match(r"heading\s+([1-6])", style_lower)
    if heading_match:
        level = int(heading_match.group(1))
        return f"{'#' * level} {text}"

    if "list" in style_lower or paragraph_has_numbering(paragraph):
        return f"- {text}"

    return text


def paragraph_has_numbering(paragraph: Paragraph) -> bool:
    ppr = paragraph._p.pPr
    return bool(ppr is not None and ppr.numPr is not None)


def word_table_to_markdown(table: Table, table_index: int) -> list[str]:
    lines = [f"## Table {table_index}"]
    for row in table.rows:
        cells = [
            clean_inline_markdown(cell.text.replace("\n", " "))
            for cell in row.cells
        ]
        cells = [cell for cell in cells if cell]
        if cells:
            lines.append("- " + " | ".join(cells))
    return lines if len(lines) > 1 else []


def read_legacy_word_document(path: Path) -> str:
    """Extract text from legacy .doc files when macOS textutil is available."""
    textutil = shutil.which("textutil")
    if not textutil:
        raise SystemExit(
            "Legacy .doc input requires macOS `textutil`. Convert the file to "
            ".docx, or run this on macOS with textutil available."
        )

    result = subprocess.run(
        [textutil, "-convert", "txt", "-stdout", str(path)],
        check=False,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
    )
    if result.returncode != 0:
        raise SystemExit(
            "Failed to extract text from legacy .doc input with textutil: "
            f"{result.stderr.strip()}"
        )
    return result.stdout


def strip_rtf(raw: str) -> str:
    """Best-effort text extraction for simple RTF notes."""
    text = re.sub(r"\\'[0-9a-fA-F]{2}", " ", raw)
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", text)
    text = text.replace("{", " ").replace("}", " ")
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def read_rtf_document(path: Path, raw: str) -> str:
    """Extract RTF text, preferring macOS textutil when available."""
    textutil = shutil.which("textutil")
    if textutil:
        result = subprocess.run(
            [textutil, "-convert", "txt", "-stdout", str(path)],
            check=False,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout

    return strip_rtf(raw)


def read_pdf_document(path: Path) -> str:
    """Extract deck-ready text from a text-based PDF."""
    if PdfReader is None:
        raise SystemExit(
            "pypdf is required for PDF inputs. Install dependencies with "
            "`pip install -r requirements.txt`. "
            f"Original import error: {PDF_IMPORT_ERROR}"
        )

    try:
        reader = PdfReader(path)
    except Exception as exc:
        raise SystemExit(f"Failed to open PDF input {path}: {exc}") from exc

    pages: list[list[str]] = []
    title_lines: list[str] = []
    for page_number, page in enumerate(reader.pages, start=1):
        try:
            page_text = page.extract_text() or ""
        except Exception as exc:
            raise SystemExit(
                f"Failed to extract text from PDF page {page_number}: {exc}"
            ) from exc

        raw_lines = clean_pdf_lines(page_text)
        if page_number == 1:
            title_lines = extract_pdf_title_lines(raw_lines)
            raw_lines = raw_lines[len(title_lines) :]

        lines = merge_pdf_lines(raw_lines)
        if page_number == 1:
            lines = trim_pdf_preamble(lines)

        if lines:
            pages.append(lines)

    if not pages:
        raise SystemExit(
            f"No extractable text found in {path}. Convert scanned pages with OCR "
            "or provide a text/Word source document."
        )

    title = " ".join(title_lines) if title_lines else path.stem.replace("-", " ").title()
    body = "\n".join(line for page in pages for line in page)
    return f"# {title}\n\n{body}"


def clean_pdf_lines(text: str) -> list[str]:
    """Repair common PDF extraction noise while preserving original line breaks."""
    text = re.sub(r"(\w)-\n(\w)", r"\1\2", text)
    lines = []
    for raw_line in text.splitlines():
        line = re.sub(r"\s+", " ", raw_line).strip()
        if line and not is_pdf_noise_line(line):
            lines.append(line)
    return lines


def is_pdf_noise_line(line: str) -> bool:
    """Filter footers, footnotes, and caption fragments that make poor slide bullets."""
    lower = line.lower()
    if re.fullmatch(r"\d{1,3}", line):
        return True
    if "arxiv:" in lower or "contact:" in lower:
        return True
    if "leaderboard:" in lower or line.startswith(("http://", "https://")):
        return True
    if re.match(r"^\d+https?://", line):
        return True
    if re.match(r"^\d+[A-Z][A-Za-z].*", line):
        return True
    if re.match(r"^[A-Za-z0-9_-]{8,}/edit$", line):
        return True
    if re.match(r"^(Figure|Table)\s+\d+:", line):
        return True
    return False


def merge_pdf_lines(lines: list[str]) -> list[str]:
    """Merge wrapped PDF text lines into paragraph-sized bullets."""
    merged: list[str] = []
    paragraph: list[str] = []

    def flush_paragraph() -> None:
        if paragraph:
            merged.append(clean_inline_markdown(" ".join(paragraph)))
            paragraph.clear()

    for line in lines:
        if is_pdf_standalone_line(line):
            flush_paragraph()
            merged.append(line)
            continue

        paragraph.append(line)
        paragraph_text = " ".join(paragraph)
        if line.endswith((".", "?", "!", ";", ":")) or len(paragraph_text) > 420:
            flush_paragraph()

    flush_paragraph()
    return merged


def is_pdf_standalone_line(line: str) -> bool:
    if is_plain_text_heading(line):
        return True
    if re.match(r"^[-*+]\s+\S+", line):
        return True
    if re.match(r"^\d+[.)]\s+\S+", line):
        return True
    return False


def trim_pdf_preamble(lines: list[str]) -> list[str]:
    """Drop title-page metadata before the first standard document heading."""
    for index, line in enumerate(lines):
        if line.lower() in COMMON_DOCUMENT_HEADINGS:
            return lines[index:]
    return lines


def extract_pdf_title_lines(first_page_lines: list[str]) -> list[str]:
    """Infer a short title from the first prominent lines of a PDF."""
    title_lines = []
    for line in first_page_lines[:4]:
        if not line[0].isalnum():
            break
        if "," in line and title_lines:
            break
        if len(line) >= 8:
            title_lines.append(clean_inline_markdown(line))
        if len(title_lines) == 2:
            break

    return title_lines


def parse_sections(text: str, fallback_title: str) -> tuple[str, list[Section]]:
    title = fallback_title
    current = Section("Overview", [])
    sections: list[Section] = []

    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            continue

        heading = re.match(r"^(#{1,4})\s+(.+)$", line)
        if heading:
            level = len(heading.group(1))
            heading_text = clean_inline_markdown(heading.group(2))
            if level == 1 and title == fallback_title:
                title = heading_text
                continue
            if current.bullets or current.title != "Overview":
                sections.append(current)
            current = Section(heading_text, [])
            continue

        if is_plain_text_heading(line):
            if current.bullets or current.title != "Overview":
                sections.append(current)
            current = Section(clean_inline_markdown(line), [])
            continue

        bullet = re.match(r"^[-*+]\s+(.+)$", line)
        numbered = re.match(r"^\d+[.)]\s+(.+)$", line)
        content = None
        if bullet:
            content = bullet.group(1)
        elif numbered:
            content = numbered.group(1)
        elif len(line) > 25:
            content = line

        if content:
            current.bullets.append(clean_inline_markdown(content))

    if current.bullets or current.title != "Overview":
        sections.append(current)

    if not sections:
        sentences = re.split(r"(?<=[.!?])\s+", text.strip())
        sections = [Section("Overview", [clean_inline_markdown(s) for s in sentences if s])]

    return title, sections


def is_plain_text_heading(line: str) -> bool:
    """Detect headings commonly produced by Word/RTF plain-text export."""
    if len(line) > 80:
        return False
    if line.lower() in COMMON_DOCUMENT_HEADINGS:
        return True
    if re.match(r"^[一二三四五六七八九十]+[、.．]\s*\S+", line):
        return True
    if re.match(r"^\d+(?:\.\d+)*\s+[A-Z][A-Za-z0-9 /&().+-]{2,}$", line):
        return True
    return False


def clean_inline_markdown(text: str) -> str:
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"\*([^*]+)\*", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"<[^>]+>", "", text)
    return re.sub(r"\s+", " ", text).strip()


def chunk_section(section: Section, max_bullets: int) -> list[Section]:
    bullets = [b for b in section.bullets if b]
    if not bullets:
        return [section]

    chunks = []
    for index in range(0, len(bullets), max_bullets):
        chunk = bullets[index : index + max_bullets]
        suffix = "" if index == 0 else f" ({index // max_bullets + 1})"
        chunks.append(Section(f"{section.title}{suffix}", chunk))
    return chunks


def select_slide_sections(sections: list[Section], max_content_slides: int) -> list[Section]:
    """Prefer broad section coverage before adding continuation slides."""
    non_empty_sections = []
    for section in sections:
        bullets = clean_bullets_for_slides(section.bullets)
        if bullets:
            non_empty_sections.append(Section(section.title, bullets))

    selected = [
        Section(section.title, section.bullets[:5])
        for section in non_empty_sections[:max_content_slides]
    ]

    if len(selected) >= max_content_slides:
        return selected

    for section in non_empty_sections:
        for continuation in chunk_section(section, 5)[1:]:
            selected.append(continuation)
            if len(selected) >= max_content_slides:
                return selected

    return selected


def clean_bullets_for_slides(bullets: list[str]) -> list[str]:
    return [bullet for bullet in bullets if is_slide_worthy_bullet(bullet)]


def is_slide_worthy_bullet(bullet: str) -> bool:
    text = bullet.strip()
    if len(text) < 25:
        return False
    if text[0].islower():
        return False

    lower = text.lower()
    if "http" in lower:
        return False
    if sum(1 for token in re.findall(r"\b\w+\b", text) if len(token) == 1) >= 8:
        return False
    broken_pdf_fragments = ("or chestr", "pr ogr", "sur f", "t ask", "t erminal")
    if any(fragment in lower for fragment in broken_pdf_fragments):
        return False
    noisy_prefixes = (
        "dataset category ",
        "figure ",
        "method gaia ",
        "spa ",
        "table ",
        "wald interval method",
        "webarena and assistantbench",
    )
    if lower.startswith(noisy_prefixes):
        return False
    if "leaderboard" in lower and ("baseline" in lower or "method" in lower):
        return False
    if re.search(r"\bLevel 1\b.*\bLevel 2\b.*\bLevel 3\b", text):
        return False
    if len(re.findall(r"\bM1\b", text)) >= 3:
        return False
    if text.count("\u00b1") >= 2 or text.count("\u2013") >= 5:
        return False
    return True


def extract_pdf_figure_assets(input_path: Path, output_path: Path) -> dict[str, FigureAsset]:
    if input_path.suffix.lower() != ".pdf":
        return {}
    if fitz is None:
        raise SystemExit(
            "PyMuPDF is required to extract PDF figures. Install dependencies with "
            "`pip install -r requirements.txt`. "
            f"Original import error: {FITZ_IMPORT_ERROR}"
        )

    asset_dir = output_path.parent / f"{output_path.stem}_assets"
    asset_dir.mkdir(parents=True, exist_ok=True)

    assets: dict[str, FigureAsset] = {}
    document = fitz.open(input_path)
    try:
        for spec in PDF_FIGURE_SPECS:
            asset = extract_pdf_figure_asset(document, asset_dir, spec)
            if asset:
                assets[asset.target_section] = asset
    finally:
        document.close()
    return assets


def extract_pdf_figure_asset(document, asset_dir: Path, spec: dict[str, object]) -> FigureAsset | None:
    label = str(spec["label"])
    for page_index in range(document.page_count):
        page = document[page_index]
        caption_rects = page.search_for(f"{label}:")
        if not caption_rects:
            continue

        caption_rect = caption_rects[0]
        crop_rect = fitz.Rect(
            50,
            float(spec["top"]),
            page.rect.width - 50,
            max(float(spec["top"]) + 40, caption_rect.y0 - 8),
        )
        output_path = asset_dir / f"{label.lower().replace(' ', '-')}.png"
        pixmap = page.get_pixmap(
            matrix=fitz.Matrix(3, 3),
            clip=crop_rect,
            alpha=False,
        )
        pixmap.save(output_path)
        return FigureAsset(
            label=label,
            target_section=str(spec["target_section"]),
            caption=str(spec["caption"]),
            path=output_path,
        )

    return None


def figure_for_section(
    section: Section,
    figures: dict[str, FigureAsset],
) -> FigureAsset | None:
    for target_section, figure in figures.items():
        if section.title.startswith(target_section):
            return figure
    return None


def build_deck(
    input_path: Path,
    output_path: Path,
    title: str | None,
    max_slides: int,
    theme_name: str,
) -> None:
    if Presentation is None:
        raise SystemExit(
            "python-pptx is required. Install dependencies with "
            "`pip install -r requirements.txt`. "
            f"Original import error: {PPTX_IMPORT_ERROR}"
        )

    theme = resolve_theme(theme_name)
    text = read_document(input_path)
    deck_title, sections = parse_sections(text, input_path.stem.replace("-", " ").title())
    if title:
        deck_title = title

    slide_sections = select_slide_sections(sections, max(1, max_slides - 2))
    figures = extract_pdf_figure_assets(input_path, output_path)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, deck_title, input_path, theme)
    for index, section in enumerate(slide_sections, start=1):
        add_content_slide(
            prs,
            section,
            index,
            len(slide_sections),
            input_path,
            theme,
            figure_for_section(section, figures),
        )
    add_close_slide(prs, deck_title, input_path, theme)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def add_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title: str, input_path: Path, theme: dict[str, RGBColor]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, theme["panel"])

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.32), Inches(7.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = theme["accent"]
    accent.line.fill.background()

    kicker = slide.shapes.add_textbox(Inches(0.85), Inches(0.8), Inches(7.5), Inches(0.4))
    set_text(kicker, "Generated from source document", 14, theme["accent_2"], bold=True)

    title_box = slide.shapes.add_textbox(Inches(0.85), Inches(1.55), Inches(9.8), Inches(2.2))
    set_text(title_box, title, 40, theme["light"], bold=True)

    subtitle = slide.shapes.add_textbox(Inches(0.9), Inches(4.25), Inches(9.2), Inches(0.8))
    set_text(
        subtitle,
        f"Source: {input_path.name}",
        17,
        RGBColor(220, 230, 238),
        bold=False,
    )

    block = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(10.25),
        Inches(5.45),
        Inches(2.15),
        Inches(0.55),
    )
    block.fill.solid()
    block.fill.fore_color.rgb = theme["accent"]
    block.line.fill.background()
    label = slide.shapes.add_textbox(Inches(10.32), Inches(5.56), Inches(2.0), Inches(0.3))
    set_text(label, "PPTX deck", 13, theme["panel"], bold=True, align=PP_ALIGN.CENTER)


def add_content_slide(
    prs,
    section: Section,
    index: int,
    total: int,
    input_path: Path,
    theme: dict[str, RGBColor],
    figure: FigureAsset | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, theme["background"])
    has_figure = figure is not None and figure.path.exists()

    rail = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5)
    )
    rail.fill.solid()
    rail.fill.fore_color.rgb = theme["accent"]
    rail.line.fill.background()

    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.65), Inches(0.55), Inches(0.62), Inches(0.62)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = theme["panel"]
    badge.line.fill.background()
    badge_label = slide.shapes.add_textbox(Inches(0.65), Inches(0.71), Inches(0.62), Inches(0.25))
    set_text(
        badge_label,
        f"{index:02d}",
        11,
        theme["light"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    title_box = slide.shapes.add_textbox(Inches(1.45), Inches(0.55), Inches(10.5), Inches(0.85))
    set_text(title_box, section.title, 30, theme["text"], bold=True)

    card_width = 5.25 if has_figure else 8.25
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.85),
        Inches(1.65),
        Inches(card_width),
        Inches(4.9),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = theme["light"]
    card.line.color.rgb = RGBColor(225, 232, 240)

    bullet_box = slide.shapes.add_textbox(
        Inches(1.18),
        Inches(1.95),
        Inches(4.62 if has_figure else 7.55),
        Inches(4.2),
    )
    text_frame = bullet_box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    bullet_limit = 4 if has_figure else 5
    for bullet in section.bullets[:bullet_limit]:
        paragraph = text_frame.add_paragraph()
        paragraph.text = shorten(bullet, width=92 if has_figure else 135, placeholder="...")
        paragraph.level = 0
        paragraph.font.size = Pt(13 if has_figure else 16)
        paragraph.font.color.rgb = theme["text"]
        paragraph.space_after = Pt(9)
    if section.bullets:
        text_frame._txBody.remove(text_frame._txBody.p_lst[0])

    if has_figure:
        add_figure_panel(slide, figure, theme)
    else:
        add_takeaway_panel(slide, section, index, total, theme)

    footer = slide.shapes.add_textbox(Inches(0.85), Inches(7.02), Inches(11.7), Inches(0.25))
    set_text(
        footer,
        f"Source: {input_path.name}",
        9,
        theme["muted"],
        align=PP_ALIGN.RIGHT,
    )


def add_takeaway_panel(
    slide,
    section: Section,
    index: int,
    total: int,
    theme: dict[str, RGBColor],
) -> None:
    callout = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.45),
        Inches(1.65),
        Inches(2.8),
        Inches(4.9),
    )
    callout.fill.solid()
    callout.fill.fore_color.rgb = theme["panel"]
    callout.line.fill.background()

    callout_title = slide.shapes.add_textbox(Inches(9.75), Inches(2.0), Inches(2.2), Inches(0.5))
    set_text(callout_title, "Takeaway", 15, theme["accent_2"], bold=True)
    takeaway = section.bullets[0] if section.bullets else section.title
    callout_body = slide.shapes.add_textbox(Inches(9.75), Inches(2.65), Inches(2.25), Inches(2.4))
    set_text(callout_body, shorten(takeaway, width=110, placeholder="..."), 18, theme["light"])

    progress = slide.shapes.add_textbox(Inches(9.75), Inches(5.75), Inches(2.2), Inches(0.32))
    set_text(progress, f"{index} of {total}", 11, RGBColor(210, 220, 230))


def add_figure_panel(slide, figure: FigureAsset, theme: dict[str, RGBColor]) -> None:
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.35),
        Inches(1.65),
        Inches(5.9),
        Inches(4.9),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = theme["light"]
    panel.line.color.rgb = RGBColor(225, 232, 240)

    add_picture_contained(slide, figure.path, 6.55, 1.85, 5.5, 3.95)

    caption = slide.shapes.add_textbox(Inches(6.62), Inches(6.0), Inches(5.35), Inches(0.35))
    set_text(caption, figure.caption, 9, theme["muted"], bold=True, align=PP_ALIGN.CENTER)


def add_picture_contained(
    slide,
    image_path: Path,
    left: float,
    top: float,
    width: float,
    height: float,
) -> None:
    from PIL import Image

    with Image.open(image_path) as image:
        image_width, image_height = image.size

    image_ratio = image_width / image_height
    box_ratio = width / height
    if image_ratio > box_ratio:
        draw_width = width
        draw_height = width / image_ratio
    else:
        draw_height = height
        draw_width = height * image_ratio

    draw_left = left + (width - draw_width) / 2
    draw_top = top + (height - draw_height) / 2
    slide.shapes.add_picture(
        str(image_path),
        Inches(draw_left),
        Inches(draw_top),
        width=Inches(draw_width),
        height=Inches(draw_height),
    )


def add_close_slide(prs, title: str, input_path: Path, theme: dict[str, RGBColor]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, theme["panel"])

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.32), Inches(7.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = theme["accent"]
    accent.line.fill.background()

    motif = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.9),
        Inches(1.15),
        Inches(2.25),
        Inches(4.7),
    )
    motif.fill.solid()
    motif.fill.fore_color.rgb = theme["accent"]
    motif.line.fill.background()

    for offset, color in enumerate((theme["panel"], theme["accent_2"], theme["light"])):
        marker = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(10.55 + (offset * 0.45)),
            Inches(3.15),
            Inches(0.3),
            Inches(0.3),
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = color
        marker.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.95), Inches(1.0), Inches(9.4), Inches(0.8))
    set_text(title_box, "Discussion", 34, theme["light"], bold=True)

    body = slide.shapes.add_textbox(Inches(1.0), Inches(2.15), Inches(8.7), Inches(1.4))
    set_text(
        body,
        f"Use this deck as a starting point for presenting {title}. Review the generated slides against the source before sharing.",
        20,
        RGBColor(224, 232, 240),
    )

    source = slide.shapes.add_textbox(Inches(1.0), Inches(5.85), Inches(8.4), Inches(0.4))
    set_text(source, f"Source document: {input_path.name}", 13, theme["accent_2"], bold=True)


def set_text(
    shape,
    text: str,
    size: int,
    color: RGBColor,
    bold: bool = False,
    align: PP_ALIGN | None = None,
) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    if align:
        paragraph.alignment = align
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def run_codex_agent(
    input_path: Path,
    output_path: Path,
    title: str | None,
    max_slides: int,
    theme_name: str,
) -> None:
    codex = shutil.which("codex")
    if not codex:
        raise SystemExit("codex CLI was not found on PATH.")

    skill_file = hermes_skill_dir() / "SKILL.md"
    if not skill_file.exists():
        raise SystemExit(
            "Hermes PowerPoint skill is not installed. Run "
            "`python tools/hermes_powerpoint.py --install-skill` first."
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prompt = build_agent_prompt(input_path, output_path, title, max_slides, theme_name, skill_file)
    command = [
        codex,
        "--ask-for-approval",
        "never",
        "exec",
        "-C",
        str(Path.cwd()),
        "-s",
        "workspace-write",
        prompt,
    ]
    subprocess.run(command, check=True)


def build_agent_prompt(
    input_path: Path,
    output_path: Path,
    title: str | None,
    max_slides: int,
    theme_name: str,
    skill_file: Path,
) -> str:
    title_line = f"Use this deck title: {title}." if title else "Infer a concise deck title."
    return f"""Use $powerpoint from NousResearch/hermes-agent to create a PowerPoint deck.

If the skill is not loaded by name in this Codex process, read and follow:
{skill_file}

Input document:
{input_path}

Output PPTX path:
{output_path}

Requirements:
- {title_line}
- Build no more than {max_slides} slides.
- Use the Hermes PowerPoint guidance for slide structure, visual variety, and QA.
- Use the {theme_name} visual direction unless the source document strongly suggests a better palette.
- Save the final editable .pptx file at exactly the output path above.
- Do not upload, email, or send the file anywhere.
- Do not search outside this repository and the installed skill directory.
- Do not install dependencies or use network access.
- For implementation, prefer this repository's deterministic wrapper:
  `python tools/hermes_powerpoint.py {input_path} --out {output_path} --max-slides {max_slides} --theme {theme_name}`
- Run a lightweight content check before finishing and report the output path.
"""


def parse_args(argv: list[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Generate a PPTX deck from a document using Hermes PowerPoint."
    )
    parser.add_argument(
        "input",
        nargs="?",
        help="Source .md, .txt, .rtf, .pdf, .docx, .docm, or legacy .doc document.",
    )
    parser.add_argument("--out", help="Output .pptx path.")
    parser.add_argument("--title", help="Optional deck title override.")
    parser.add_argument("--max-slides", type=int, default=10)
    parser.add_argument("--theme", choices=sorted(THEMES), default="executive")
    parser.add_argument(
        "--agent",
        action="store_true",
        help="Launch codex exec and ask it to use the Hermes PowerPoint skill.",
    )
    parser.add_argument(
        "--install-skill",
        action="store_true",
        help="Install NousResearch/hermes-agent PowerPoint skill into CODEX_HOME.",
    )
    return parser.parse_args(argv)


def main(argv: list[str]) -> int:
    args = parse_args(argv)
    if args.install_skill:
        install_hermes_skill()
        if not args.input:
            return 0

    if not args.input:
        raise SystemExit("Provide an input document or use --install-skill.")

    input_path = Path(args.input).expanduser().resolve()
    if not input_path.exists():
        raise SystemExit(f"Input document not found: {input_path}")

    output_path = (
        Path(args.out).expanduser().resolve()
        if args.out
        else Path("outputs/presentations") / f"{input_path.stem}.pptx"
    )
    output_path = output_path.resolve()

    if args.agent:
        run_codex_agent(input_path, output_path, args.title, args.max_slides, args.theme)
    else:
        build_deck(input_path, output_path, args.title, args.max_slides, args.theme)
        print(f"Wrote {output_path}")

    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
